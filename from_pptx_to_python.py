"""
pptx_to_py.py  --  convert a .pptx into a Python script that regenerates it.

Usage:
    python pptx_to_py.py  IHX_Drawings_to_convert_to_py.pptx

Writes IHX_Drawings_to_convert_to_py.py next to it. The generated script
rebuilds the deck faithfully by re-injecting each shape's XML (this preserves
freeforms, dashes, arrowheads, fonts and colours, which a textbox/connector
reconstruction would drop).

Note: embedded media (pictures, charts) reference relationship IDs that are
NOT carried over by raw-XML injection. Decks built purely from drawn shapes
and text (like the IHX drawings) round-trip cleanly; picture-heavy decks won't.
"""
import os
import sys
from pptx import Presentation


def convert(in_path: str) -> str:
    prs = Presentation(in_path)
    base = os.path.splitext(in_path)[0]
    out_path = base + ".py"
    out_name = os.path.basename(base) + "_roundtrip.pptx"

    blocks = []
    for slide in prs.slides:
        shape_xmls = [shape._element.xml for shape in slide.shapes]
        blocks.append(shape_xmls)

    with open(out_path, "w", encoding="utf-8") as f:
        f.write('"""Auto-generated from %s by pptx_to_py.py."""\n' % os.path.basename(in_path))
        f.write("from pptx import Presentation\n")
        f.write("from pptx.util import Emu\n")
        f.write("from pptx.oxml import parse_xml\n\n")
        f.write("prs = Presentation()\n")
        f.write("prs.slide_width  = Emu(%d)\n" % int(prs.slide_width))
        f.write("prs.slide_height = Emu(%d)\n" % int(prs.slide_height))
        f.write("BLANK = prs.slide_layouts[6]\n\n")
        f.write("SLIDES = [\n")
        for shape_xmls in blocks:
            f.write("    [\n")
            for xml in shape_xmls:
                f.write("        %r,\n" % xml)
            f.write("    ],\n")
        f.write("]\n\n")
        f.write("for shape_xmls in SLIDES:\n")
        f.write("    slide = prs.slides.add_slide(BLANK)\n")
        f.write("    spTree = slide.shapes._spTree\n")
        f.write("    for xml in shape_xmls:\n")
        f.write("        spTree.append(parse_xml(xml))\n\n")
        f.write('prs.save(%r)\n' % out_name)
        f.write('print("saved", %r)\n' % out_name)

    return out_path


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("usage: python pptx_to_py.py <file.pptx>")
        sys.exit(1)
    p = convert(sys.argv[1])
    print("wrote", p)