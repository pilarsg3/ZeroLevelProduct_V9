"""
IHX parameter drawings (2 slides):
  1. General elevation: XZ section of the whole unit (two plenums, tube bundle,
     windowed bundle shell, central pipe, outlet riser, lateral pipe) with a
     vertical break through the long tube bundle.
  2. Central-pipe routing detail: on-axis riser -> 90 deg bend inside the upper
     plenum -> horizontal exit through the +X wall, in the true plane of the bend.
"""
import math
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

IN = 914400

# ---------------------------------------------------------------- palette
INK      = RGBColor(0x33, 0x3D, 0x4F)
STEEL    = RGBColor(0xC9, 0xD6, 0xE3)
STEEL_D  = RGBColor(0x47, 0x55, 0x69)
STEEL_L  = RGBColor(0xED, 0xF1, 0xF6)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
WALLHATCH= RGBColor(0x9A, 0xA5, 0xB1)
PLATE    = RGBColor(0xAE, 0xC3, 0xD6)

C_RAD = RGBColor(0x0F, 0x76, 0x6E)   # radius      (green)
C_HGT = RGBColor(0x1F, 0x3A, 0x5F)   # height/axial(navy)
C_LEN = RGBColor(0xC2, 0x57, 0x0C)   # length      (orange)
C_THK = RGBColor(0xB9, 0x1C, 0x1C)   # thickness   (red)
C_ANG = RGBColor(0x6D, 0x28, 0xD9)   # angle       (purple)

prs = Presentation()
prs.slide_width  = Emu(int(13.333 * IN))
prs.slide_height = Emu(int(7.5 * IN))
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- helpers
def e(v): return int(round(v * IN))

def add_text(s, x, y, w, h, text, size=12, color=INK, bold=False,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             italic=False, font="Calibri"):
    tb = s.shapes.add_textbox(e(x), e(y), e(w), e(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    for i, ln in enumerate(text.split("\n")):
        para = p if i == 0 else tf.add_paragraph()
        para.alignment = align
        r = para.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = font; r.font.color.rgb = color
    return tb

def _set_dash(ln, dash):
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))

def _set_arrows(ln, head, tail):
    if head:
        ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if tail:
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))

def line(s, x1, y1, x2, y2, color=INK, w=1.25, dash=None, head=False, tail=False):
    cxn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, e(x1), e(y1), e(x2), e(y2))
    cxn.line.color.rgb = color; cxn.line.width = Pt(w)
    ln = cxn.line._get_or_add_ln()
    _set_dash(ln, dash); _set_arrows(ln, head, tail)
    return cxn

def dim(s, x1, y1, x2, y2, color, label, lsize=11, loff=(0, 0), lalign=PP_ALIGN.CENTER):
    line(s, x1, y1, x2, y2, color=color, w=1.1, head=True, tail=True)
    mx, my = (x1 + x2) / 2 + loff[0], (y1 + y2) / 2 + loff[1]
    if label:
        add_text(s, mx - 0.85, my - 0.13, 1.7, 0.26, label, lsize, color, bold=True, align=lalign)

def freeform(s, pts, fill=None, edge=INK, w=1.4, close=True, dash=None):
    fb = s.shapes.build_freeform(e(pts[0][0]), e(pts[0][1]), scale=1.0)
    fb.add_line_segments([(e(px), e(py)) for px, py in pts[1:]], close=close)
    sh = fb.convert_to_shape()
    if fill is None: sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if edge is None: sh.line.fill.background()
    else:
        sh.line.color.rgb = edge; sh.line.width = Pt(w)
        _set_dash(sh.line._get_or_add_ln(), dash)
    return sh

def rect(s, x, y, w, h, fill=None, edge=INK, lw=1.3, dash=None):
    return freeform(s, [(x, y), (x + w, y), (x + w, y + h), (x, y + h)],
                    fill=fill, edge=edge, w=lw, dash=dash)

def polyline(s, pts, color=INK, w=1.4, dash=None):
    for i in range(len(pts) - 1):
        line(s, pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1], color=color, w=w, dash=dash)

def legend(s, x, y, items=None):
    items = items or [(C_RAD, "Radius"), (C_HGT, "Height / axial"),
                      (C_LEN, "Length"), (C_THK, "Thickness"), (C_ANG, "Angle")]
    add_text(s, x, y - 0.30, 2.4, 0.26, "Parameter colour key", 11, INK, bold=True, align=PP_ALIGN.LEFT)
    for i, (c, lab) in enumerate(items):
        yy = y + i * 0.27
        rect(s, x, yy, 0.20, 0.18, fill=c, edge=None)
        add_text(s, x + 0.30, yy - 0.04, 2.1, 0.26, lab, 10.5, INK, align=PP_ALIGN.LEFT)

# ================================================================ demo values
lp_ir, lp_wall, lp_h, lp_dr = 300., 20., 200., 320.
up_ir, up_wall, up_h, up_dr = 300., 20., 300., 320.
up_or = up_ir + up_wall
bh = 2000.
rings = [(150., 12., 14.), (220., 10., 12.), (260., 10., 12.)]   # (pitch, ir, or)
cp_ir, cp_or, cp_bend, cp_z, cp_horiz = 60., 70., 80., 100., 400.
rs_ir, rs_or, rs_h = 80., 90., 300.
lat_ir, lat_or, lat_len, lat_z = 40., 48., 300., 150.
bs_ir, bs_or, bs_wall, bs_nbars, bs_barw = 285., 300., 15., 8, 20.
bs_winh = 0.2 * bh                      # window_fraction 0.2
bs_dz   = bs_winh / 2.0                 # default gap

z_lp_top = lp_h                 # 200
z_up_bot = z_lp_top + bh        # 2200
z_up_top = z_up_bot + up_h      # 2500
z_dome_t = z_up_top + up_dr     # 2820  (riser base)
z_cp_bend = z_up_bot + cp_z     # 2300
z_cp_hor  = z_cp_bend + cp_bend # 2380
x_cp_far  = cp_bend + up_or + cp_horiz   # 800
z_clip    = z_up_top + math.sqrt(up_dr**2 - rs_or**2)   # ~2807
z_rs_top  = z_dome_t + rs_h      # 3120
z_lat     = z_dome_t + lat_z     # 2970

# windows (z ranges)
win_up = (z_up_bot - bs_dz - bs_winh, z_up_bot - bs_dz)   # (1600, 2000)
win_lo = (z_lp_top + bs_dz, z_lp_top + bs_dz + bs_winh)   # (400, 800)

# ================================================================ SLIDE 1
s1 = prs.slides.add_slide(BLANK)
add_text(s1, 0.5, 0.28, 12.3, 0.55, "Intermediate heat exchanger \u2014 general elevation",
         25, INK, bold=True, align=PP_ALIGN.LEFT, font="Georgia")
add_text(s1, 0.5, 0.84, 12.3, 0.30,
         "create_ihx  \u00b7  XZ section: lower plenum, tube bundle + windowed bundle shell, upper plenum, central pipe, outlet riser, lateral pipe  \u00b7  bundle broken to scale",
         11.5, STEEL_D, align=PP_ALIGN.LEFT, italic=True)

AX = 3.55                   # axis (model x=0) slide-x
hs = vs = 0.00205           # in / model-unit (uniform)
def MX(x):  return AX + x * hs
yDB = 6.88                  # slide-y of dome bottom (z=-lp_dr)
def BY(z):  return yDB - (z + lp_dr) * vs           # bottom band  z in [-lp_dr, 650]
zB_top = 650.
y_brk_b = BY(zB_top)
y_brk_t = y_brk_b - 0.50
def TY(z):  return y_brk_t - (z - 1750.) * vs        # top band z in [1750, z_rs_top]
zT_bot = 1750.

def mrect(my, x0, z0, x1, z1, **kw):
    rect(s1, MX(x0), my(z1), (x1 - x0) * hs, (z1 - z0) * vs, **kw)

def dome(my, zc, r, wall, lower, fill=STEEL):
    a0, a1 = (180, 360) if lower else (0, 180)
    outer = [(MX(r*math.cos(math.radians(a))), my(zc + r*math.sin(math.radians(a))))
             for a in range(a0, a1 + 1, 3)]
    inner = [(MX((r-wall)*math.cos(math.radians(a))), my(zc + (r-wall)*math.sin(math.radians(a))))
             for a in range(a1, a0 - 1, -3)]
    freeform(s1, outer + inner, fill=fill, edge=STEEL_D, w=1.4)

# ---- LOWER PLENUM -------------------------------------------------------
dome(BY, 0.0, lp_dr, lp_wall, lower=True)                      # bottom dome shell
mrect(BY, -lp_ir-lp_wall, 0., -lp_ir, z_lp_top, fill=STEEL, edge=STEEL_D, lw=1.4)  # L wall
mrect(BY,  lp_ir, 0.,  lp_ir+lp_wall, z_lp_top, fill=STEEL, edge=STEEL_D, lw=1.4)  # R wall
mrect(BY, -lp_ir, z_lp_top-lp_wall, lp_ir, z_lp_top, fill=PLATE, edge=STEEL_D, lw=1.3)  # tube sheet

# ---- UPPER PLENUM -------------------------------------------------------
mrect(TY, -up_ir, z_up_bot, up_ir, z_up_bot+up_wall, fill=PLATE, edge=STEEL_D, lw=1.3)  # bottom plate
mrect(TY, -up_ir-up_wall, z_up_bot, -up_ir, z_up_top, fill=STEEL, edge=STEEL_D, lw=1.4)
mrect(TY,  up_ir, z_up_bot,  up_ir+up_wall, z_up_top, fill=STEEL, edge=STEEL_D, lw=1.4)
dome(TY, z_up_top, up_dr, up_wall, lower=False)                # top dome shell

# ---- BUNDLE SHELL (windowed) : outer line w/ gaps at window rows --------
def shell_wall(side):
    xo = side * bs_or
    # bottom band: z_lp_top .. zB_top, skip window_lo (400..800)
    segs_b = [(z_lp_top, win_lo[0]), (win_lo[1], zB_top)]
    for z0, z1 in segs_b:
        if z1 > z0: line(s1, MX(xo), BY(z0), MX(xo), BY(z1), color=STEEL_D, w=1.5)
    # top band: zT_bot .. z_up_bot, skip window_up (1600..2000)
    segs_t = [(zT_bot, z_up_bot)]
    out = []
    for z0, z1 in segs_t:
        # remove the window_up gap
        if win_up[0] > z0: out.append((z0, min(win_up[0], z1)))
        if win_up[1] < z1: out.append((max(win_up[1], z0), z1))
    for z0, z1 in out:
        if z1 > z0: line(s1, MX(xo), TY(z0), MX(xo), TY(z1), color=STEEL_D, w=1.5)
for sd in (+1, -1):
    shell_wall(sd)
add_text(s1, MX(bs_or)+0.06, TY(win_up[1])-0.02, 1.8, 0.4,
         "bundle shell\n(windowed)", 9.5, STEEL_D, italic=True, align=PP_ALIGN.LEFT)

# ---- TUBES (representative, both sides, both bands) ---------------------
def tubes(my, z0, z1):
    for pitch, tir, tor in rings:
        for sd in (+1, -1):
            xc = sd * pitch
            rect(s1, MX(xc-tor), my(z1), 2*tor*hs, (z1-z0)*vs, fill=STEEL, edge=STEEL_D, lw=0.6)
            line(s1, MX(xc), my(z0), MX(xc), my(z1), color=WHITE, w=0.7)
tubes(BY, z_lp_top, zB_top)
tubes(TY, zT_bot, z_up_bot)

# ---- CENTRAL PIPE -------------------------------------------------------
# vertical (bottom band): z ~ z_lp_top-12 .. zB_top
mrect(BY, -cp_or, z_lp_top-12, cp_or, zB_top, fill=STEEL, edge=STEEL_D, lw=1.0)
line(s1, MX(0), BY(z_lp_top-12), MX(0), BY(zB_top), color=WHITE, w=1.2)
# vertical (top band): zT_bot .. z_cp_bend
mrect(TY, -cp_or, zT_bot, cp_or, z_cp_bend, fill=STEEL, edge=STEEL_D, lw=1.0)
line(s1, MX(0), TY(zT_bot), MX(0), TY(z_cp_bend), color=WHITE, w=1.2)
# bend + horizontal (top band): build outer band along centreline
def cp_centreline():
    pts = [(0., z_cp_bend)]
    for a in range(180, 89, -3):            # arc centre (cp_bend, z_cp_bend), 180->90
        pts.append((cp_bend + cp_bend*math.cos(math.radians(a)),
                    z_cp_bend + cp_bend*math.sin(math.radians(a))))
    pts.append((x_cp_far, z_cp_hor))
    return pts
def band(cl, hw):
    m=len(cl); up=[]; lo=[]
    for i,(x,z) in enumerate(cl):
        if i==0: tx,tz=cl[1][0]-x, cl[1][1]-z
        elif i==m-1: tx,tz=x-cl[i-1][0], z-cl[i-1][1]
        else: tx,tz=cl[i+1][0]-cl[i-1][0], cl[i+1][1]-cl[i-1][1]
        L=math.hypot(tx,tz) or 1.; nx,nz=-tz/L, tx/L
        up.append((x+hw*nx, z+hw*nz)); lo.append((x-hw*nx, z-hw*nz))
    return up+lo[::-1]
clc = cp_centreline()
freeform(s1, [(MX(x), TY(z)) for x,z in band(clc, cp_or)], fill=STEEL, edge=STEEL_D, w=1.0)
freeform(s1, [(MX(x), TY(z)) for x,z in band(clc, cp_ir)], fill=WHITE, edge=STEEL_D, w=0.6)
add_text(s1, MX(440), TY(z_cp_hor)+0.20, 2.7, 0.24,
         "central pipe (see detail)", 9.5, STEEL_D, italic=True, align=PP_ALIGN.LEFT)
line(s1, MX(560), TY(z_cp_hor)-0.05, MX(560), TY(z_cp_hor)+0.18,
     color=STEEL_D, w=0.6)

# ---- OUTLET RISER + cap -------------------------------------------------
mrect(TY, -rs_or, z_clip, -rs_ir, z_rs_top, fill=STEEL, edge=STEEL_D, lw=1.2)
mrect(TY,  rs_ir, z_clip,  rs_or, z_rs_top, fill=STEEL, edge=STEEL_D, lw=1.2)
mrect(TY, -rs_or, z_rs_top-10, rs_or, z_rs_top, fill=STEEL, edge=STEEL_D, lw=1.2)  # top cap

# ---- LATERAL PIPE (off riser +X) ---------------------------------------
mrect(TY, rs_ir, z_lat-lat_or, rs_or+lat_len, z_lat+lat_or, fill=STEEL, edge=STEEL_D, lw=1.0)
line(s1, MX(rs_ir), TY(z_lat), MX(rs_or+lat_len), TY(z_lat), color=WHITE, w=1.0)

# ---- centreline + break zigzags ----------------------------------------
line(s1, MX(0), BY(-lp_dr)-0.05, MX(0), BY(z_lp_top-12), color=INK, w=0.7, dash="dashDot")
line(s1, MX(0), TY(z_cp_bend), MX(0), TY(z_rs_top)+0.05, color=INK, w=0.7, dash="dashDot")
def zigzag(y, x0, x1, amp=0.06, segs=14):
    pts=[]
    for i in range(segs+1):
        x=x0+(x1-x0)*i/segs; dy = amp if i%2 else -amp
        if i in (0,segs): dy=0
        pts.append((x, y+dy))
    polyline(s1, pts, color=STEEL_D, w=1.2)
zigzag(y_brk_b, MX(-lp_ir), MX(lp_ir))
zigzag(y_brk_t, MX(-up_ir), MX(up_ir))

# ================================ dimensions (slide 1) ===================
# --- left stack: heights (navy) ---
# bundle_height (far left, broken)
xbh = MX(-up_ir-up_wall) - 1.35
line(s1, MX(-bs_or), BY(z_lp_top), xbh-0.05, BY(z_lp_top), color=C_HGT, w=0.6)
line(s1, MX(-up_ir-up_wall), TY(z_up_bot), xbh-0.05, TY(z_up_bot), color=C_HGT, w=0.6)
line(s1, xbh, BY(z_lp_top), xbh, y_brk_b+0.05, color=C_HGT, w=1.1, head=True)
line(s1, xbh, TY(z_up_bot), xbh, y_brk_t-0.05, color=C_HGT, w=1.1, head=True)
polyline(s1, [(xbh-0.06,y_brk_b+0.03),(xbh+0.06,y_brk_b-0.03)], color=C_HGT, w=1.2)
polyline(s1, [(xbh-0.06,y_brk_t+0.03),(xbh+0.06,y_brk_t-0.03)], color=C_HGT, w=1.2)
add_text(s1, xbh-0.30, (TY(z_up_bot)+BY(z_lp_top))/2-0.12, 1.5, 0.24, "bundle_height",
         9.5, C_HGT, bold=True, align=PP_ALIGN.RIGHT)
# lower_plenum_height
xlh = MX(-lp_ir-lp_wall) - 0.55
line(s1, MX(-lp_ir-lp_wall), BY(0.), xlh-0.05, BY(0.), color=C_HGT, w=0.6)
line(s1, MX(-lp_ir-lp_wall), BY(z_lp_top), xlh-0.05, BY(z_lp_top), color=C_HGT, w=0.6)
dim(s1, xlh, BY(0.), xlh, BY(z_lp_top), C_HGT, "", lsize=9)
add_text(s1, xlh-1.95, (BY(0.)+BY(z_lp_top))/2-0.12, 1.85, 0.24, "lower_plenum_height",
         9.5, C_HGT, bold=True, align=PP_ALIGN.RIGHT)
# upper_plenum_height
xuh = MX(-up_ir-up_wall) - 0.55
line(s1, MX(-up_ir-up_wall), TY(z_up_bot), xuh-0.05, TY(z_up_bot), color=C_HGT, w=0.6)
line(s1, MX(-up_ir-up_wall), TY(z_up_top), xuh-0.05, TY(z_up_top), color=C_HGT, w=0.6)
dim(s1, xuh, TY(z_up_bot), xuh, TY(z_up_top), C_HGT, "", lsize=9)
add_text(s1, xuh-1.95, (TY(z_up_bot)+TY(z_up_top))/2-0.12, 1.85, 0.24, "upper_plenum_height",
         9.5, C_HGT, bold=True, align=PP_ALIGN.RIGHT)

# --- inner radii (green) ---
line(s1, MX(0), BY(70.), MX(lp_ir), BY(70.), color=C_RAD, w=1.0, head=True, tail=True)
add_text(s1, MX(150)-0.85, BY(70.)+0.05, 1.7, 0.22, "lower_plenum_inner_radius",
         8.5, C_RAD, bold=True, align=PP_ALIGN.CENTER)
line(s1, MX(0), TY(2440.), MX(up_ir), TY(2440.), color=C_RAD, w=1.0, head=True, tail=True)
add_text(s1, MX(150)-0.85, TY(2440.)-0.26, 1.7, 0.22, "upper_plenum_inner_radius",
         8.5, C_RAD, bold=True, align=PP_ALIGN.CENTER)

# --- dome radii (green, diagonal) ---
line(s1, MX(0), BY(0.), MX(lp_dr*math.cos(math.radians(225))),
     BY(lp_dr*math.sin(math.radians(225))), color=C_RAD, w=1.0, tail=True)
add_text(s1, MX(-226)-1.5, BY(-226.)-0.20, 1.5, 0.40, "lower_plenum_\ndome_radius",
         8.5, C_RAD, bold=True, align=PP_ALIGN.RIGHT)
line(s1, MX(0), TY(z_up_top), MX(up_dr*math.cos(math.radians(125))),
     TY(z_up_top+up_dr*math.sin(math.radians(125))), color=C_RAD, w=1.0, tail=True)
add_text(s1, MX(-184)-1.5, TY(2762.)+0.06, 1.5, 0.40, "upper_plenum_\ndome_radius",
         8.5, C_RAD, bold=True, align=PP_ALIGN.RIGHT)

# --- wall leaders (red) ---
line(s1, MX(-lp_ir-lp_wall/2), BY(40.), MX(-lp_ir-lp_wall)-0.45, BY(40.)+0.30, color=C_THK, w=0.7)
add_text(s1, MX(-lp_ir-lp_wall)-1.80, BY(40.)+0.20, 1.35, 0.22, "lower_plenum_wall",
         8.5, C_THK, bold=True, align=PP_ALIGN.RIGHT)
line(s1, MX(up_ir+up_wall/2), TY(2470.), MX(up_ir+up_wall)+0.55, TY(2470.)-0.12, color=C_THK, w=0.7)
add_text(s1, MX(up_ir+up_wall)+0.60, TY(2470.)-0.24, 1.4, 0.22, "upper_plenum_wall",
         8.5, C_THK, bold=True, align=PP_ALIGN.LEFT)

# --- riser_height (navy, left of riser) ---
xrh = MX(-rs_or) - 0.30
line(s1, MX(-rs_or), TY(z_dome_t), xrh+0.05, TY(z_dome_t), color=C_HGT, w=0.6)
line(s1, MX(-rs_or), TY(z_rs_top), xrh+0.05, TY(z_rs_top), color=C_HGT, w=0.6)
dim(s1, xrh, TY(z_dome_t), xrh, TY(z_rs_top), C_HGT, "", lsize=9)
add_text(s1, xrh-1.55, (TY(z_dome_t)+TY(z_rs_top))/2-0.12, 1.45, 0.24, "riser_height",
         9, C_HGT, bold=True, align=PP_ALIGN.RIGHT)

# --- lateral_pipe_length (orange, above lateral) ---
yL = TY(z_lat+lat_or) - 0.18
line(s1, MX(rs_or), TY(z_lat+lat_or), MX(rs_or), yL+0.04, color=C_LEN, w=0.6)
line(s1, MX(rs_or+lat_len), TY(z_lat+lat_or), MX(rs_or+lat_len), yL+0.04, color=C_LEN, w=0.6)
dim(s1, MX(rs_or), yL, MX(rs_or+lat_len), yL, C_LEN, "lateral_pipe_length",
    lsize=8.5, loff=(0,-0.16))

# --- lateral_pipe_z_offset (navy, right of lateral) ---
xlz = MX(rs_or+lat_len) + 0.30
line(s1, MX(rs_or+lat_len), TY(z_dome_t), xlz+0.05, TY(z_dome_t), color=C_HGT, w=0.6)
line(s1, MX(rs_or+lat_len), TY(z_lat), xlz+0.05, TY(z_lat), color=C_HGT, w=0.6)
dim(s1, xlz, TY(z_dome_t), xlz, TY(z_lat), C_HGT, "", lsize=9)
add_text(s1, xlz+0.10, (TY(z_dome_t)+TY(z_lat))/2-0.12, 1.45, 0.24, "lateral_pipe_\nz_offset",
         8.5, C_HGT, bold=True, align=PP_ALIGN.LEFT)

legend(s1, 11.25, 5.95)

# ================================================================ SLIDE 2
s2 = prs.slides.add_slide(BLANK)
add_text(s2, 0.5, 0.28, 12.3, 0.55, "IHX central pipe \u2014 routing detail",
         25, INK, bold=True, align=PP_ALIGN.LEFT, font="Georgia")
add_text(s2, 0.5, 0.84, 12.3, 0.30,
         "On-axis vertical riser \u2192 90\u00b0 bend inside the upper plenum \u2192 horizontal exit through the +X wall (shown in the XZ plane of the bend)",
         11.5, STEEL_D, align=PP_ALIGN.LEFT, italic=True)

# local frame: model (x,z) with z measured from the upper-plenum bottom plate (z_up_bot=0)
S = 0.0072
OX, OZ = 3.30, 5.00
def QX(x): return OX + x * S
def QZ(z): return OZ - z * S

# upper-plenum bottom plate + wall context (hatched)
plate_x0, plate_x1 = -90., up_ir          # plate drawn from near axis out to inner wall
rect(s2, QX(plate_x0), QZ(0.+up_wall), (plate_x1-plate_x0)*S, up_wall*S,
     fill=STEEL_L, edge=STEEL_D, lw=1.0)
add_text(s2, QX(-90.)-0.02, QZ(0.)+0.06, 2.2, 0.22, "upper-plenum bottom plate",
         9, STEEL_D, italic=True, align=PP_ALIGN.LEFT)
# +X cylindrical wall (the pipe pierces it)
rect(s2, QX(up_ir), QZ(up_h), up_wall*S, (up_h-(-30.))*S, fill=STEEL_L, edge=STEEL_D, lw=1.0)
add_text(s2, QX(up_or)+0.04, QZ(up_h)+0.10, 1.7, 0.4, "upper-plenum\n+X wall", 9, STEEL_D,
         italic=True, align=PP_ALIGN.LEFT)

# central-pipe path (local): vertical from below plate -> 90 arc -> horizontal exit
def cl_local():
    pts = [(0., -20.)]
    pts.append((0., cp_z))                         # up to bend start (bend centre z = cp_z)
    for a in range(180, 89, -3):
        pts.append((cp_bend + cp_bend*math.cos(math.radians(a)),
                    cp_z + cp_bend*math.sin(math.radians(a))))
    pts.append((x_cp_far, cp_z + cp_bend))
    return pts
clL = cl_local()
freeform(s2, [(QX(x), QZ(z)) for x,z in band(clL, cp_or)], fill=STEEL, edge=STEEL_D, w=1.6)
freeform(s2, [(QX(x), QZ(z)) for x,z in band(clL, cp_ir)], fill=WHITE, edge=STEEL_D, w=0.9)
# centreline
polyline(s2, [(QX(x), QZ(z)) for x,z in clL], color=INK, w=1.0, dash="dashDot")

Cx, Cz = cp_bend, cp_z      # bend centre (local)
# bend centre mark
line(s2, QX(Cx)-0.07, QZ(Cz), QX(Cx)+0.07, QZ(Cz), color=C_RAD, w=1.0)
line(s2, QX(Cx), QZ(Cz)-0.07, QX(Cx), QZ(Cz)+0.07, color=C_RAD, w=1.0)

# ---- dimensions (slide 2) ----
# central_pipe_bend_radius (green): C to arc midpoint
amid = math.radians(135)
Pm = (Cx + cp_bend*math.cos(amid), Cz + cp_bend*math.sin(amid))
line(s2, QX(Cx), QZ(Cz), QX(Pm[0]), QZ(Pm[1]), color=C_RAD, w=1.1, tail=True)
add_text(s2, QX(Pm[0])-1.95, QZ(Pm[1])-0.34, 1.85, 0.24, "central_pipe_bend_radius",
         9.5, C_RAD, bold=True, align=PP_ALIGN.RIGHT)
add_text(s2, QX(Cx)+0.06, QZ(Cz)+0.04, 0.5, 0.22, "C", 11, C_RAD, bold=True)
# 90 deg angle arc (purple)
ang = [(QX(Cx + 0.55*cp_bend*math.cos(math.radians(a))),
        QZ(Cz + 0.55*cp_bend*math.sin(math.radians(a)))) for a in range(90, 181, 5)]
polyline(s2, ang, color=C_ANG, w=1.3)
line(s2, QX(Cx), QZ(Cz), QX(Cx), QZ(Cz+0.72*cp_bend), color=C_ANG, w=0.6, dash="dash")
line(s2, QX(Cx), QZ(Cz), QX(Cx-0.72*cp_bend), QZ(Cz), color=C_ANG, w=0.6, dash="dash")
add_text(s2, QX(Cx-0.95*cp_bend)-0.10, QZ(Cz+0.62*cp_bend)-0.12, 1.0, 0.22, "90\u00b0",
         10, C_ANG, bold=True, align=PP_ALIGN.CENTER)
# central_pipe_z_offset (navy): plate (z=0) to bend centre (z=cp_z)
xzo = QX(-cp_or) - 0.55
line(s2, QX(-cp_or), QZ(0.), xzo-0.05, QZ(0.), color=C_HGT, w=0.6)
line(s2, QX(Cx), QZ(Cz), xzo-0.05, QZ(Cz), color=C_HGT, w=0.6)
dim(s2, xzo, QZ(0.), xzo, QZ(Cz), C_HGT, "", lsize=9)
add_text(s2, xzo-1.85, (QZ(0.)+QZ(Cz))/2-0.12, 1.75, 0.24, "central_pipe_z_offset",
         9.5, C_HGT, bold=True, align=PP_ALIGN.RIGHT)
# central_pipe_horiz_len (orange): from (cp_bend+up_or) to far end  [equals cp_horiz]
zexit = cp_z + cp_bend
x_hl0 = cp_bend + up_or                       # 400: where horiz_len begins
yhl = QZ(zexit + cp_or) - 0.40
line(s2, QX(x_hl0), QZ(zexit+cp_or), QX(x_hl0), yhl+0.05, color=C_LEN, w=0.6, dash="dash")
line(s2, QX(x_cp_far), QZ(zexit+cp_or), QX(x_cp_far), yhl+0.05, color=C_LEN, w=0.6)
dim(s2, QX(x_hl0), yhl, QX(x_cp_far), yhl, C_LEN, "central_pipe_horiz_len", lsize=9.5, loff=(0,-0.16))
# central_pipe_inner_radius (green) + wall (red) at the open end
endx = x_cp_far
line(s2, QX(endx), QZ(zexit), QX(endx), QZ(zexit+cp_or), color=C_RAD, w=0.9, head=True, tail=True)
add_text(s2, QX(endx)-0.55, QZ(zexit+cp_or)-0.28, 1.9, 0.22, "central_pipe_inner_radius",
         8.5, C_RAD, bold=True, align=PP_ALIGN.LEFT)
line(s2, QX(endx), QZ(zexit-cp_or), QX(endx)+0.5, QZ(zexit-cp_or)+0.28, color=C_THK, w=0.6)
line(s2, QX(endx), QZ(zexit-cp_ir), QX(endx)+0.5, QZ(zexit-cp_ir)+0.28, color=C_THK, w=0.6)
dim(s2, QX(endx)+0.45, QZ(zexit-cp_or)+0.28, QX(endx)+0.45, QZ(zexit-cp_ir)+0.28, C_THK, "", lsize=9)
add_text(s2, QX(endx)+0.58, QZ(zexit-(cp_or+cp_ir)/2)+0.10, 1.5, 0.22, "central_pipe_wall",
         8.5, C_THK, bold=True, align=PP_ALIGN.LEFT)

legend(s2, 11.25, 1.55)

prs.save("IHX_Drawings.pptx")
print("saved")